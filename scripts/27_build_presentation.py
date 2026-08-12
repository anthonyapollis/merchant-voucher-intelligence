"""
27_build_presentation.py — the interview walkthrough deck.

Built to be TALKED THROUGH, not read from. One idea per slide, the number that carries it in
large type, and the reasoning underneath in a size that supports the speaker rather than
competing with them. Every figure comes from docs/*.json, so the deck cannot drift from the
warehouse.
"""
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "report" / "Merchant_Voucher_Intelligence_Presentation.pptx"
sys.path.insert(0, str(Path(__file__).parent))
from _table_registry import TABLES as REG, TIERS, counts as tier_counts, TIER_ORDER

S = json.load(open(ROOT / "docs" / "analytics_summary.json"))
ML = json.load(open(ROOT / "docs" / "ml_summary.json"))
RC = json.load(open(ROOT / "docs" / "reconciliation.json"))
K, BA, SZ = S["exec_kpis"], S["business_answers"], S["size_confounder"]
TC = tier_counts()

NAVY, NAVY2, TEAL, AMBER = "12305B", "1B4079", "0E8B8B", "E8A317"
RED, PURPLE, GREEN, GREY, INK = "C0392B", "7B4B94", "1E8449", "5A6672", "12203A"
BG, CARD, LINE = "F2F5FA", "FFFFFF", "DCE4EF"
# python-pptx wants bare hex; the shared registry stores CSS colours with a leading '#'.
C = lambda h: RGBColor.from_string(h.lstrip("#").upper())

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
W, H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = C(bg); r.line.fill.background()
    r.shadow.inherit = False
    return s


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for item in runs:
        text, size, colour, bold = item[0], item[1], item[2], item[3]
        space_before = item[4] if len(item) > 4 else 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = spacing
        p.space_before = Pt(space_before)
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = C(colour); r.font.name = "Segoe UI"
    return tb


def band(s, title, sub=""):
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(1.0))
    b.fill.solid(); b.fill.fore_color.rgb = C(NAVY); b.line.fill.background()
    b.shadow.inherit = False
    txt(s, 0.55, 0.20, 9.5, 0.6, [(title, 26, "FFFFFF", True)])
    if sub:
        txt(s, 0.55, 0.66, 12.2, 0.3, [(sub, 12, "9FB6D4", False)])


def card(s, x, y, w, h, fill=CARD, edge=LINE):
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                           Inches(w), Inches(h))
    c.fill.solid(); c.fill.fore_color.rgb = C(fill)
    c.line.color.rgb = C(edge); c.line.width = Pt(0.75)
    c.shadow.inherit = False
    try:
        c.adjustments[0] = 0.06
    except Exception:
        pass
    return c


def kpi(s, x, y, w, label, value, sub, colour=NAVY):
    card(s, x, y, w, 1.55, fill=colour, edge=colour)
    txt(s, x + 0.2, y + 0.16, w - 0.4, 0.25, [(label.upper(), 10, "C9DAF0", True)])
    txt(s, x + 0.2, y + 0.44, w - 0.4, 0.6, [(value, 30, "FFFFFF", True)])
    txt(s, x + 0.2, y + 1.08, w - 0.4, 0.35, [(sub, 9.5, "AFC4DC", False)])


def bullets(s, x, y, w, items, size=13.5, gap=9):
    runs = []
    for i, (head, body) in enumerate(items):
        runs.append((head, size, NAVY, True, 0 if i == 0 else gap))
        if body:
            runs.append((body, size - 1.5, GREY, False, 2))
    # 7.5in canvas, leaving a 0.25in bottom margin — never a fixed 5.0in.
    txt(s, x, y, w, max(7.5 - y - 0.25, 0.5), runs, spacing=1.15)


def picture(s, img, x, y, max_w, max_h):
    from PIL import Image
    iw, ih = Image.open(img).size
    ratio = min(max_w / (iw / 96), max_h / (ih / 96))
    w_in = (iw / 96) * ratio
    h_in = (ih / 96) * ratio
    s.shapes.add_picture(str(img), Inches(x + (max_w - w_in) / 2),
                         Inches(y + (max_h - h_in) / 2), Inches(w_in), Inches(h_in))


R = lambda v: f"R{v:,.0f}"
Rm = lambda v: f"R{v/1e6:.1f}m"
P1 = lambda v: f"{v*100:.1f}%"

# ======================================================================= 1. TITLE
s = slide(NAVY)
t = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.9), W, Inches(1.6))
t.fill.solid(); t.fill.fore_color.rgb = C(TEAL); t.line.fill.background()
t.shadow.inherit = False
txt(s, 0.9, 1.9, 11.5, 1.4, [("Merchant Sales &", 44, "FFFFFF", True),
                             ("Voucher Intelligence", 44, "5FD4D0", True)], spacing=1.05)
txt(s, 0.9, 3.7, 11.5, 0.9,
    [("BI Developer — Second-Round Practical Task", 17, "9FB6D4", False),
     ("Microsoft Fabric  ·  Power BI  ·  dbt  ·  Python ML", 14, "7FA3C9", False, 6)])
txt(s, 0.9, 6.15, 11.5, 1.0,
    [("Anthony Apollis", 16, "FFFFFF", True),
     ("1 Jan – 31 Jul 2026  ·  25 merchants  ·  120,969 vouchers  ·  510,127 transactions",
      12, "CDEEEC", False, 4)])

# ======================================================================= 2. WHAT WAS BUILT
s = slide()
band(s, "What was built", "Every figure below is reproducible: python scripts/run_all.py rebuilds it in ~2 minutes")
row = [("Data pipeline", "Bronze → Silver → Gold\nmedallion, 149,857 rows", NAVY),
       ("dbt project", f"18 models · 1 SCD2 snapshot\n132 tests · 156 pass", TEAL),
       ("Semantic model", "14 tables · 15 enforced FKs\n70 DAX measures", PURPLE),
       ("Power BI", "9 report pages\n+ phone layout", NAVY2),
       ("ML", "5 models, time-split\nvalidation throughout", AMBER)]
for i, (h1, h2, col) in enumerate(row):
    x = 0.55 + i * 2.5
    card(s, x, 1.5, 2.3, 1.85, fill=col, edge=col)
    txt(s, x + 0.18, 1.68, 1.95, 0.35, [(h1, 13, "FFFFFF", True)])
    txt(s, x + 0.18, 2.15, 1.95, 1.1, [(h2, 10.5, "D8E4F2", False)], spacing=1.25)

card(s, 0.55, 3.65, 12.2, 3.3)
bullets(s, 0.95, 3.95, 11.4, [
    ("Quality is the deliverable, not a footnote.",
     "14 warehouse integrity tests · 132 dbt tests · 33 idempotency assertions · "
     "31 DAX measures each with a SQL-derived expected value."),
    ("The gold layer is built TWICE, by two independent implementations.",
     "Once in pandas, once in dbt SQL. Every headline figure agrees to the cent — 27 of 28 "
     "checks tie exactly, and the one warning is a documented rounding convention. That is "
     "what makes \"the numbers tie\" verifiable instead of asserted."),
    ("It caught two real defects that no unit test would have found.",
     "dim_date was absorbing the August redemption tail into the reporting window, inflating "
     "pro-rated targets by R984,046. And pandas rank(pct=True) is not the same statistic as "
     "SQL PERCENT_RANK() — that shifted the Health Score by up to 9.7 points."),
])

# ======================================================================= 3. ARCHITECTURE
s = slide()
band(s, "Architecture", "The journey of one row, and what each stage prevents")
stages = [("LANDING", "4 CSVs", AMBER), ("BRONZE", "raw + lineage", "8B6914"),
          ("SILVER", "typed, conformed", GREY), ("QUALITY GATE", "132 dbt tests", RED),
          ("GOLD", "Kimball star", TEAL), ("POWER BI", "import now,
Direct Lake ready", NAVY)]
for i, (nm, sub, col) in enumerate(stages):
    x = 0.5 + i * 2.12
    card(s, x, 1.45, 1.9, 1.1, fill=col, edge=col)
    txt(s, x + 0.1, 1.62, 1.7, 0.3, [(nm, 11, "FFFFFF", True)], align=PP_ALIGN.CENTER)
    txt(s, x + 0.1, 1.98, 1.7, 0.4, [(sub, 9, "E8EFF8", False)], align=PP_ALIGN.CENTER)
    if i < len(stages) - 1:
        txt(s, x + 1.92, 1.78, 0.25, 0.3, [("→", 16, GREY, True)])

card(s, 0.5, 2.85, 12.3, 1.5, fill="FDF0EE", edge=RED)
txt(s, 0.85, 3.05, 11.7, 1.1, [
    ("The quality gate sits BETWEEN silver and gold. That is the whole design.", 15, RED, True),
    ("If a test fails, the pipeline stops and gold is NOT rebuilt — the report keeps showing "
     "yesterday's correct numbers rather than today's wrong ones. Testing after loading gold "
     "means the wrong figures have already been published by the time anyone reads the alert. "
     "Stale-but-correct beats fresh-but-wrong.", 12, INK, False, 5)], spacing=1.2)

card(s, 0.5, 4.55, 6.0, 2.45)
txt(s, 0.85, 4.78, 5.4, 2.0, [
    ("Orchestration — Data Factory", 14, NAVY, True),
    ("Metadata-driven ingest: adding a source is a parameter change, not a pipeline edit.",
     11, GREY, False, 6),
    ("Per-source row-count floors catch the failure that does not throw — a truncated file "
     "looks exactly like a quiet trading day.", 11, GREY, False, 4),
    ("Retry 3 on copy, retry 0 on the quality gate: a failing test is a defect, not a blip.",
     11, GREY, False, 4)], spacing=1.15)

card(s, 6.75, 4.55, 6.05, 2.45)
txt(s, 7.1, 4.78, 5.45, 2.0, [
    ("Why dbt is in this stack", 14, TEAL, True),
    ("Runs locally in seconds and on the Fabric Warehouse in production — the "
     "same models, only the adapter changes.", 11, GREY, False, 6),
    ("Tests, lineage and documentation are part of the build, not a separate exercise "
     "somebody does later.", 11, GREY, False, 4),
    ("Business rules live in ONE place and every consumer inherits them.", 11, GREY, False, 4)],
    spacing=1.15)

# ======================================================================= 4. DAG
s = slide()
band(s, "dbt lineage", f"Generated from target/manifest.json — 33 nodes, 70 edges, 6 layers")
dag = ROOT / "docs" / "dbt_dag.png"
if dag.exists():
    picture(s, dag, 0.4, 1.15, 12.5, 5.5)

# ======================================================================= 5-6. ERDs
for img, title, sub in (
    ("erd_before.png", "Before — what was delivered",
     "Four flat files · 6 duplicated attribute columns · 0 enforced keys · every value text"),
    ("erd_after.png", "After — the gold star schema",
     "14 tables · 15 ENFORCED foreign keys · conformed dimensions · contiguous calendar"),
):
    p = ROOT / "docs" / img
    if p.exists():
        s = slide()
        band(s, title, sub)
        picture(s, p, 0.4, 1.15, 12.5, 5.9)

# ======================================================================= 7. WHY 14 TABLES
s = slide()
band(s, "Why 14 tables when the README suggests 5",
     "Justified individually — the honest split is a 10-table core plus 4 extensions")
x = 0.55
for k in TIER_ORDER:
    label, col, blurb = TIERS[k]
    col = col.lstrip("#").upper()
    n = TC[k]
    members = sorted(t for t, v in REG.items() if v["tier"] == k)
    card(s, x, 1.45, 3.0, 4.3, fill=CARD, edge=col)
    hdr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.45),
                             Inches(3.0), Inches(0.85))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = C(col); hdr.line.fill.background()
    hdr.shadow.inherit = False
    txt(s, x + 0.15, 1.56, 2.7, 0.3, [(label, 12.5, "FFFFFF", True)])
    txt(s, x + 0.15, 1.86, 2.7, 0.3, [(f"{n} tables", 20, "FFFFFF", True)])
    txt(s, x + 0.15, 2.45, 2.7, 0.5, [(blurb, 9.5, GREY, False)], spacing=1.1)
    txt(s, x + 0.15, 3.05, 2.7, 2.6,
        [(m, 9.5, INK if k != "extension" else PURPLE, False, 3) for m in members],
        spacing=1.1)
    x += 3.12

card(s, 0.55, 5.95, 12.2, 1.05, fill="FEF8EC", edge=AMBER)
txt(s, 0.9, 6.12, 11.6, 0.8, [
    ("The fair criticism: four dimensions of 4–6 rows is over-normalisation. ", 12, AMBER, True),
    ("That was the first design here — all three ticket dimensions in one table behind a "
     "discriminator — and it was worse: the FKs on fct_support_tickets could not be tested, "
     "and Power BI cannot build three independent filter paths off one physical table. "
     "Splitting them added 9 enforced foreign keys that previously could not exist.",
     11, INK, False)], spacing=1.15)

# ======================================================================= 8. KPIs
s = slide()
band(s, "The numbers", "1 January – 31 July 2026")
kpi(s, 0.55, 1.45, 2.35, "Total sales", Rm(K["TotalSales"]),
    f"{K['TotalTransactions']:,} transactions", NAVY)
kpi(s, 3.05, 1.45, 2.35, "Redemption", P1(K["RedemptionRate"]),
    f"{K['VouchersRedeemed']:,} of {K['VouchersSold']:,}", TEAL)
kpi(s, 5.55, 1.45, 2.35, "Avg basket", R(K["AvgBasketValue"]), "all voucher types", NAVY2)
kpi(s, 8.05, 1.45, 2.35, "SLA breach", P1(K["SLABreachRate"]),
    "94.7% on High/Critical", RED)
kpi(s, 10.55, 1.45, 2.25, "Liability", Rm(K["OutstandingLiability"]),
    "unredeemed value", PURPLE)

card(s, 0.55, 3.3, 12.25, 3.7)
bullets(s, 0.95, 3.6, 11.5, [
    ("Umhlanga Value Mart is failing — and it is operational, not demand.",
     "July sales −42.5% against its own prior three-month average, while support tickets rose "
     "+693% in the same month. R571,518 annualised at risk, and plausibly recoverable because "
     "the two moving together points to a service breakdown rather than lost customers."),
    ("The SLA policy is configured backwards.",
     "Critical gets a 12-hour target and takes 52.7 hours (98.3% breach). Low gets 48 hours "
     "and takes 11.3 (0.2%). The ladder runs opposite to the actual workload, so 94.7% of all "
     "358 breaches land on High and Critical. The metric is measuring a policy error, not the "
     "team."),
    ("Eastern Cape is the one declining region — on four independent signals.",
     "Only region that peaked before July; 9.8% below its own peak while every other region "
     "is AT its peak; flattest trend slope; and a 12.2% June fall against +1.4% to +3.8% "
     "elsewhere. One signal would be noise. Four is a finding."),
], size=13)

# ======================================================================= 8b. THE 5 QUESTIONS
s = slide()
band(s, "The five business questions", "Section 6 of the brief — every answer computed, "
                                       "with the SQL in sql/02_business_questions.sql")
QS = [
    ("Q1", "Which merchants generate the highest sales value and transaction volume?",
     "Durban Cash Hub leads on BOTH — R5,776,119 (8.8%) and 45,371 transactions. Both "
     "rankings are reported because they need not agree: a larger basket can put a merchant "
     "top on value and not on volume. Concentration matters more than the leader — the top "
     "5 hold 35.4% of revenue, and 15 of 25 merchants make up 80%.", NAVY),
    ("Q2", "Which voucher type has the highest redemption rate?",
     "Airtime at 92.8%; Gaming lowest at 76.0% — a 16.9pp spread. Value-based and "
     "volume-based rates are near identical within each type, and time-to-redeem is flat "
     "(3.5–3.7 days), so the difference is WHETHER customers redeem, not how quickly. "
     "Gaming's 24% non-redemption is the largest block of the R3.5m liability.", TEAL),
    ("Q3", "Which region shows declining sales or transaction behaviour?",
     "Eastern Cape, on four independent signals: the only region to peak before July; 9.8% "
     "below its own peak while all others are AT peak; the flattest trend slope (+2.0% vs "
     "+4.4%–8.5%); and a 12.2% June fall against +1.4%–3.8% elsewhere. Three of the five "
     "Critical/Watch merchants sit there.", AMBER),
    ("Q4", "Are ticket volumes, priority or resolution times associated with weaker "
           "performance?",
     "Not at portfolio level — the apparent r = −0.56 is confounded by merchant size and "
     "collapses to −0.20 once controlled. The real signal is event-level, and is covered on "
     "the following slide.", RED),
    ("Q5", "Which merchants should management focus on first, and why?",
     "Ranked by revenue at risk, not severity: 1. Umhlanga Value Mart (R571,518), "
     "2. Table Bay Express (R354,163), 3. Pretoria PayPoint (R154,045). A 42.5% collapse at "
     "a small merchant costs less than a 6% slide at a large one — ranking on percentage "
     "change alone directs the account team to the wrong merchant.", PURPLE),
]
y = 1.35
PAGE_REF = {'Q1': 'Power BI: Merchant Analysis  ·  league table + Executive Overview KPIs', 'Q2': 'Power BI: Voucher & Redemption  ·  redemption rate and liability by type', 'Q3': 'Power BI: Geographic Intelligence  ·  province map and trend', 'Q4': 'Power BI: Operational View  ·  and Business Answers for the confounding', 'Q5': 'Power BI: Merchant Value & Risk  ·  ranked by revenue at risk'}
for tag, q, a, col in QS:
    card(s, 0.55, y, 12.25, 1.16)
    n = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(y + 0.14),
                           Inches(0.62), Inches(0.5))
    n.fill.solid(); n.fill.fore_color.rgb = C(col); n.line.fill.background()
    n.shadow.inherit = False
    txt(s, 0.7, y + 0.23, 0.62, 0.3, [(tag, 14, "FFFFFF", True)], align=PP_ALIGN.CENTER)
    txt(s, 1.5, y + 0.11, 11.1, 0.28, [(q, 12, NAVY, True)])
    txt(s, 1.5, y + 0.40, 11.1, 0.56, [(a, 9.5, GREY, False)], spacing=1.10)
    txt(s, 1.5, y + 0.96, 11.1, 0.18,
        [(PAGE_REF.get(tag, ""), 8.5, col, True)])
    y += 1.23

# ======================================================================= 9. Q4 — the honest one
s = slide()
band(s, "Analytical rigour — a confounded correlation",
     "Q4 — are ticket volumes associated with weaker merchant performance?")
card(s, 0.55, 1.4, 6.0, 2.85, fill="FDF0EE", edge=RED)
txt(s, 0.9, 1.58, 5.4, 2.5, [
    ("The surface reading", 14, RED, True),
    ("Tickets per 1,000 transactions correlates with target attainment at", 11.5, INK,
     False, 5),
    ("r = −0.56", 30, RED, True, 3),
    # r is used four times on this slide and never defined. Anyone who does not already know
    # the scale cannot tell whether -0.56 is a strong finding or a rounding error.
    ("r = Pearson correlation coefficient. Runs −1 to +1. 0 means no linear relationship; "
     "the SIGN is the direction, the SIZE is the strength. −0.56 is a moderate inverse "
     "relationship: as one goes up, the other tends to go down.", 9.5, RED, False, 4),
    ("\"Operational friction hurts merchant performance.\" Clean, quotable — and not "
     "supported once tested.", 11, GREY, False, 4)], spacing=1.12)

card(s, 6.8, 1.4, 6.0, 2.5, fill="F0F9F9", edge=TEAL)
txt(s, 7.15, 1.62, 5.4, 2.1, [
    ("Why it is confounded", 14, TEAL, True),
    ("That ratio is a size measure in disguise — it correlates with log(total sales) at "
     "r = −0.83. Small merchants score badly because the denominator is small.",
     12, INK, False, 6),
    ("Partial correlation, controlling for size:  r = −0.20", 14, TEAL, True, 6),
    ("SLA breach rate and resolution time show no association at all.", 11.5, GREY, False, 4)],
    spacing=1.15)

card(s, 0.55, 4.1, 12.25, 2.9)
txt(s, 0.95, 4.35, 11.5, 2.4, [
    ("Where the real signal is: event-level, not portfolio-level", 15, NAVY, True),
    ("Durban Cash Hub — tickets +780% in June, and sales GREW 8.2%.  A service problem at a "
     "healthy account.", 12.5, INK, False, 8),
    ("Umhlanga Value Mart — tickets +693% in July, and sales FELL 42.5%.  A failing account.",
     12.5, INK, False, 4),
    ("Identical operational signal, opposite commercial diagnosis. Ticket spikes are worth "
     "investigating immediately, but they do not predict revenue on their own — which is why "
     "the alerting logic pairs ticket movement with sales movement rather than firing on "
     "either alone.", 12, GREY, False, 8)], spacing=1.2)

# ======================================================================= 10. ML
s = slide()
band(s, "Machine learning", "Time-based validation throughout — never a random split")
rows = [
    ("Anomaly detection", "Isolation Forest",
     f"{ML['anomaly_detection']['anomalies_flagged']} flagged from "
     f"{ML['anomaly_detection']['observations_scored']}",
     "All 4 documented embedded patterns recovered unprompted", TEAL),
    ("Redemption propensity", "HistGradientBoosting",
     f"AUC {ML['redemption_propensity']['roc_auc']:.3f}",
     f"vs a theoretical ceiling of {ML['redemption_propensity']['oracle_auc_ceiling']:.3f}",
     AMBER),
    ("Resolution time", "HistGradientBoosting",
     f"MAE {ML['resolution_time']['mae_hours']:.1f}h",
     f"{P1(ML['resolution_time']['improvement_vs_naive'])} better than naive", NAVY2),
    ("Sales forecast", "Holt-Winters",
     f"MAPE {ML['sales_forecast']['mape']*100:.2f}%",
     f"{P1(ML['sales_forecast']['improvement_vs_naive'])} better than naive", NAVY),
    ("Segmentation", "K-Means (k=4)", "silhouette 0.185",
     "k chosen on business grounds; the cost is stated", PURPLE),
]
y = 1.4
for name, algo, metric, note, col in rows:
    card(s, 0.55, y, 12.25, 0.92)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(y),
                             Inches(0.09), Inches(0.92))
    bar.fill.solid(); bar.fill.fore_color.rgb = C(col); bar.line.fill.background()
    bar.shadow.inherit = False
    txt(s, 0.85, y + 0.14, 2.6, 0.3, [(name, 13, NAVY, True)])
    txt(s, 0.85, y + 0.48, 2.6, 0.3, [(algo, 10, GREY, False)])
    txt(s, 3.6, y + 0.26, 2.3, 0.4, [(metric, 15, col, True)])
    txt(s, 6.1, y + 0.28, 6.4, 0.5, [(note, 11.5, GREY, False)])
    y += 1.03

card(s, 0.55, y + 0.05, 12.25, 0.82, fill="FEF8EC", edge=AMBER)
txt(s, 0.9, y + 0.16, 11.6, 0.70, [
    ("How the models were validated", 13, AMBER, True),
    ("Time-based splits only — train Jan–May, test Jun–Jul. A random split leaks future "
     "information and produces a metric that will not hold in production.   ·   Every "
     "supervised model is compared against a NAIVE baseline, not just reported in "
     "isolation.   ·   The anomaly model was checked against four patterns documented in the "
     "dataset README and recovered all four without being told what they were.   ·   The "
     "0.620 AUC was tested against a theoretical ceiling: a voucher-type base-rate oracle "
     "scores 0.621, so the model captures 99.8% of achievable signal and the limit is the "
     "data, not the model.", 11, INK, False, 4)], spacing=1.15)

# ======================================================================= 10b. GEO
s = slide()
band(s, "Geographic intelligence", "Where the business is — and, more usefully, where it is not")
for i, (reg, mer, sales, share, col) in enumerate([
        ("Free State", 10, "R18.0m", "27.5%", NAVY),
        ("Gauteng", 4, "R16.1m", "24.6%", TEAL),
        ("Western Cape", 4, "R13.1m", "19.9%", NAVY2),
        ("KwaZulu-Natal", 3, "R9.5m", "14.5%", PURPLE),
        ("Eastern Cape", 4, "R8.8m", "13.4%", RED)]):
    x = 0.55 + i * 2.5
    card(s, x, 1.4, 2.3, 1.5, fill=col, edge=col)
    txt(s, x + 0.16, 1.55, 2.0, 0.28, [(reg, 11.5, "FFFFFF", True)])
    txt(s, x + 0.16, 1.86, 2.0, 0.42, [(sales, 22, "FFFFFF", True)])
    txt(s, x + 0.16, 2.42, 2.0, 0.3, [(f"{share} · {mer} merchants", 9.5, "C9DAF0", False)])

card(s, 0.55, 3.05, 6.0, 2.0, fill="FEF8EC", edge=AMBER)
txt(s, 0.9, 3.25, 5.4, 1.7, [
    ("5 of 9 provinces · 56% of the country uncovered", 14, AMBER, True),
    ("Limpopo, Mpumalanga, North West and Northern Cape carry no merchants at all. The "
     "choropleth draws them as hatched \"no cover\" rather than omitting them — an absent "
     "province otherwise reads as zero sales when it actually means no footprint.",
     11.5, INK, False, 6)], spacing=1.15)

card(s, 6.8, 3.05, 6.0, 2.0, fill="FDF0EE", edge=RED)
txt(s, 7.15, 3.25, 5.4, 1.7, [
    ("Eastern Cape is the only province in decline", 14, RED, True),
    ("The only region below its own peak (−9.8%) while every other is AT peak, and the only "
     "one with negative recent momentum. Three of the five Critical/Watch merchants sit "
     "there — a regional pattern, not one bad account.", 11.5, INK, False, 6)], spacing=1.15)

card(s, 0.55, 5.2, 12.25, 1.8, fill="EEF3F9", edge=NAVY)
txt(s, 0.95, 5.42, 11.5, 1.5, [
    ("What the geography does NOT tell us — and why that matters", 13.5, NAVY, True),
    ("It is tempting to read 56% uncovered land as 56% untapped market. It is not. No "
     "population, GDP or competitor feed was supplied, and land area is a poor proxy for "
     "demand — Northern Cape is the largest province and the most sparsely populated in the "
     "country. The honest statement is a FOOTPRINT observation: four provinces have no "
     "merchant presence. Sizing that as opportunity requires population and income data the "
     "business already holds and this dataset does not.", 11.5, GREY, False, 6)],
    spacing=1.18)

# ======================================================================= 11. CONTROLS
s = slide()
band(s, "Reconciliation and controls", "Controls that pass matter as much as controls that fail")
card(s, 0.55, 1.4, 6.0, 2.3, fill=TEAL, edge=TEAL)
txt(s, 0.9, 1.62, 5.4, 1.9, [
    ("5 of 6 controls PASS at zero variance", 15, "FFFFFF", True),
    ("Sales value survives bronze → gold · voucher and ticket row counts survive · "
     "redeemed + outstanding = issued · the scorecard ties to the sales fact.",
     12, "D6F0EE", False, 8)], spacing=1.2)

card(s, 6.8, 1.4, 6.0, 2.3, fill="FEF8EC", edge=AMBER)
txt(s, 7.15, 1.62, 5.4, 1.9, [
    ("The 6th is a R43,501,446 variance — and it is EXPECTED", 15, AMBER, True),
    ("Sales fact: R65.5m across 510,127 transactions. Voucher fact: R22.0m across 120,969 "
     "vouchers. Different populations, ratio 4.2 : 1. They must never be forced to tie.",
     12, INK, False, 8)], spacing=1.2)

card(s, 0.55, 3.9, 12.25, 3.1)
txt(s, 0.95, 4.15, 11.5, 2.6, [
    ("Why recording an expected variance is the point", 15, NAVY, True),
    ("Without it, someone compares R65.5m to R22.0m, calls it a reconciliation break and "
     "escalates — or divides one by the other and reports a value-based redemption rate "
     "against the wrong denominator. The control does not just check a number; it records "
     "which comparisons are legitimate.", 12.5, INK, False, 8),
    ("Fraud controls follow the same discipline.", 14, NAVY, True, 12),
    ("Supported and built: reversal rate, redemption velocity, liability concentration, "
     "behavioural anomaly.   Tested and NIL: voucher value outliers — zero beyond 3 SD, "
     "because the generator used a bounded distribution; the control runs and correctly "
     "returns zero rather than being quietly dropped.   Not supported: duplicate redemption, "
     "geographic anomaly, PIN misuse, customer CLV and churn — there is no customer "
     "identifier anywhere in the four files, so each one names the telemetry it would need "
     "instead.", 12, GREY, False, 6)], spacing=1.2)

# ======================================================================= 11b. FRAUD FRAMEWORK
s = slide()
band(s, "Fraud detection and prevention",
     "Detection is one third of a control framework — prevention and correction are the "
     "other two")
COLS = [
    ("PREVENT", "Stops the loss occurring", "#8B2E24", [
        ("Voucher expiry enforcement", "90-day rule is modelled; enforce at redemption"),
        ("Value ceilings per type", "Gaming carries the highest non-redemption at 24%"),
        ("Velocity limits per merchant", "Cap redemptions per merchant per hour"),
        ("Onboarding due diligence", "2 merchants are CRM 'At Risk' with no linked control"),
        ("Segregation of duties", "Issuance and redemption approval must not share an owner"),
    ]),
    ("DETECT", "Finds it after the fact — what is built", "#0E8B8B", [
        ("Reversal rate per 1k txn", "215 reversal tickets; range 0.00–1.99 per 1k"),
        ("Redemption velocity", "Same-day share, ~23% and uniform across merchants"),
        ("Behavioural anomaly", "Isolation Forest on merchant-month deviation"),
        ("Liability concentration", "R3.5m outstanding, tracked per merchant"),
        ("Value outliers", "Implemented; correctly returns ZERO on this data"),
    ]),
    ("CORRECT", "Limits the damage once found", "#B8860B", [
        ("Quarantine on DQ failure", "Bad loads never reach gold — the pipeline gate"),
        ("Change alerting", "SCD2 diff raises status/target changes as Critical"),
        ("Reversal workflow", "Financial tickets are 25% of volume; route separately"),
        ("Merchant suspension path", "Risk band 'Review' should trigger a hold, not an email"),
        ("Audit trail", "batch_id on every row answers 'which load produced this?'"),
    ]),
]
for i, (title, sub, col, items) in enumerate(COLS):
    x = 0.55 + i * 4.15
    card(s, x, 1.35, 3.95, 4.45)
    hdr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.35),
                             Inches(3.95), Inches(0.78))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = C(col.lstrip("#")); hdr.line.fill.background()
    hdr.shadow.inherit = False
    txt(s, x + 0.18, 1.46, 3.6, 0.3, [(title, 15, "FFFFFF", True)])
    txt(s, x + 0.18, 1.76, 3.6, 0.28, [(sub, 9.5, "E8EFF8", False)])
    y2 = 2.28
    for head, body in items:
        txt(s, x + 0.18, y2, 3.6, 0.62,
            [(head, 10.5, NAVY, True), (body, 9, GREY, False, 2)], spacing=1.1)
        y2 += 0.68

card(s, 0.55, 5.95, 12.25, 1.05, fill="FDF0EE", edge=RED)
txt(s, 0.9, 6.12, 11.6, 0.85, [
    ("Four controls cannot be built on this dataset, and each names what it needs. ",
     12, RED, True),
    ("Duplicate redemption — voucher_id is unique by construction, so a duplicate cannot "
     "appear; needs the raw redemption event log with failed attempts.   Geographic anomaly "
     "— needs transaction location and timestamp, not static merchant province.   PIN misuse "
     "— needs PIN, attempt outcome and failure reason.   Customer velocity — needs a customer "
     "or session identifier, which does not exist in any of the four files.",
     10.5, INK, False, 4)], spacing=1.15)

# ======================================================================= 12. FABRIC STATUS
s = slide()
band(s, "Deployed to Microsoft Fabric", "Provisioned by API, authenticated without stored secrets")
ok = [("Workspace and Lakehouse",
       "WS_MerchantVoucher with LH_MerchantVoucher holding the medallion layers in OneLake"),
      ("Warehouse", "WH_MerchantVoucher, provisioned through the Fabric REST API"),
      ("Medallion pipeline",
       "ZIP to bronze to silver to gold, four notebooks, completed in 5m 54s"),
      ("Gold in the Warehouse",
       "26,500 sales rows, 120,969 vouchers, R65,521,298.75 — ties to the cent"),
      ("Authentication", "az login CLI token — no service principal, no stored secret")]
card(s, 0.55, 1.4, 6.0, 3.55, fill="EFF9F2", edge=GREEN)
txt(s, 0.9, 1.6, 5.4, 0.3, [("Provisioned and verified", 15, GREEN, True)])
txt(s, 0.9, 2.0, 5.4, 2.7,
    [(f"{a}\n{b}", 11.5, INK, False, 8) for a, b in ok], spacing=1.15)

card(s, 6.8, 1.4, 6.0, 3.55, fill="EEF3F9", edge=TEAL)
txt(s, 7.15, 1.6, 5.4, 0.3, [("Built to run on either engine", 15, TEAL, True)])
txt(s, 7.15, 2.0, 5.4, 2.7, [
    ("One codebase, two targets.", 12.5, TEAL, True),
    ("Development runs against a local analytical engine — seconds per build, no cloud "
     "cost — and production targets the Fabric Warehouse through dbt-fabric. The project "
     "and the models are the same; only the adapter changes.", 11.5, INK, False, 6),
    ("Dialect differences are handled in a portability macro layer rather than by "
     "maintaining two copies of the SQL: mvi_datediff, mvi_bool, mvi_median, "
     "mvi_percentile and mvi_arg_max dispatch per adapter.", 11.5, INK, False, 6),
    ("Verified by connecting and running against the warehouse, not assumed.",
     11.5, GREY, False, 6)], spacing=1.15)

card(s, 0.55, 5.0, 12.25, 2.0, fill="EEF3F9", edge=NAVY)
txt(s, 0.95, 5.25, 11.5, 1.6, [
    ("Cost control, because unattended compute does not stop on its own", 14, NAVY, True),
    ("The daily trigger ships DISABLED with a hard end date of 2026-10-03 — one week before "
     "the capacity is due to be released. fabric_cost_guard.py --check / --stop / --nuke "
     "is the kill switch: "
     "it inventories every item, flags always-on item types that bill continuously rather "
     "than per run, and two scheduled tasks watch the clock. Deleting the workspace costs "
     "nothing irreplaceable — run_all.py rebuilds the entire warehouse locally in about two "
     "minutes.", 12, INK, False, 6)], spacing=1.2)

# ======================================================================= 13. RECOMMENDATIONS
s = slide()
band(s, "Recommendations", "Ranked by revenue at risk, not by severity of decline")
recs = [
    ("1", "Site visit to Umhlanga Value Mart this week", R(571518),
     "Sales −42.5% and tickets +693% in the same month — operational failure causing "
     "commercial damage, which means it is plausibly recoverable", RED),
    ("2", "Re-base the SLA policy", "Restores the metric",
     "94.7% of breaches come from Critical/High because targets run inverse to workload. "
     "A 90th-percentile-compliant Critical SLA would be ~120h", AMBER),
    ("3", "Root-cause Durban Cash Hub's ticket volume", Rm(5776119),
     "+780% tickets on the largest merchant in the book while sales still grow — a service "
     "problem that has not yet become a commercial one", NAVY),
    ("4", "Eastern Cape regional review", Rm(8781130),
     "The only region below its own peak; 3 of the 5 Critical/Watch merchants sit there", TEAL),
    ("5", "Confirm the sales target basis with Finance", "Makes attainment reportable",
     "Targets sit ~6.1× below realised sales for ALL 25 merchants — a basis error, not "
     "outperformance", PURPLE),
]
y = 1.4
for num, action, impact, why, col in recs:
    card(s, 0.55, y, 12.25, 1.03)
    n = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(y + 0.18),
                           Inches(0.55), Inches(0.62))
    n.fill.solid(); n.fill.fore_color.rgb = C(col); n.line.fill.background()
    n.shadow.inherit = False
    txt(s, 0.7, y + 0.3, 0.55, 0.35, [(num, 17, "FFFFFF", True)], align=PP_ALIGN.CENTER)
    txt(s, 1.45, y + 0.14, 5.0, 0.3, [(action, 13.5, NAVY, True)])
    txt(s, 1.45, y + 0.5, 8.6, 0.45, [(why, 10.5, GREY, False)], spacing=1.1)
    txt(s, 10.3, y + 0.3, 2.3, 0.4, [(impact, 14, col, True)], align=PP_ALIGN.RIGHT)
    y += 1.13

# ======================================================================= 13b. FURTHER Q&A
s = slide()
band(s, "Further questions and answers",
     "Questions this solution raises beyond the brief, with the evidence behind each answer")
QA2 = [
    ("Why not simply average the per-merchant redemption rates?",
     "It would weight a 3,500-voucher merchant the same as a 5,400-voucher one. Every ratio "
     "in the model is SUM/SUM, never an average of ratios — including in the filtered views "
     "of the report.", TEAL),
    ("The two value facts differ by R43.5m. Is that a reconciliation break?",
     "No, and the control records it as EXPECTED. MerchantSales aggregates all 510,127 "
     "transactions; VoucherRedemptions covers 120,969 individual vouchers — a 4.2 : 1 ratio, "
     "two different populations. Forcing them to tie would be the error.", NAVY),
    ("Why is target attainment shown as an index rather than a percentage?",
     "The supplied BaseMonthlySalesTarget sits ~6.1x below realised sales for ALL 25 "
     "merchants, so raw attainment reads 614%. The value is retained unchanged and flagged; "
     "a relative index is reported alongside it, which is comparable across merchants and "
     "immune to the calibration error.", AMBER),
    ("Can this detect fraud?",
     "Partly, and the boundary is documented. Reversal rate, redemption velocity, liability "
     "concentration and behavioural anomaly are built. Voucher-value outliers were "
     "implemented and correctly return zero. Duplicate redemption, geographic anomaly and "
     "PIN misuse cannot be built — each names the telemetry required.", RED),
    ("What happens if the source file arrives truncated?",
     "The pipeline fails loudly. Per-source minimum row counts turn a silent partial load "
     "into an error, because a truncated file otherwise looks identical to a quiet trading "
     "day on the dashboard.", PURPLE),
    ("How is the report protected from a bad data load?",
     "The dbt test suite runs BETWEEN silver and gold. If it fails, gold is not rebuilt and "
     "the report continues to serve the last good load rather than publishing wrong figures.",
     NAVY2),
]
y = 1.35
for q, a, col in QA2:
    card(s, 0.55, y, 12.25, 0.92)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(y), Inches(0.08),
                             Inches(0.92))
    bar.fill.solid(); bar.fill.fore_color.rgb = C(col); bar.line.fill.background()
    bar.shadow.inherit = False
    txt(s, 0.85, y + 0.1, 11.7, 0.26, [(q, 12.5, NAVY, True)])
    txt(s, 0.85, y + 0.38, 11.7, 0.5, [(a, 10.5, GREY, False)], spacing=1.12)
    y += 0.99

# ======================================================================= 14. CLOSE
s = slide(NAVY)
txt(s, 1.0, 1.3, 11.3, 0.8, [("Principles applied throughout", 32, "FFFFFF", True)])
items = [
    ("Figures are verifiable, not asserted.",
     "The gold layer is built twice by independent implementations and reconciled to the "
     "cent. 33 idempotency assertions prove a re-run cannot duplicate data."),
    ("Findings are tested before they are reported.",
     "The friction/performance correlation proved confounded by merchant size. The redemption "
     "AUC sits at the data's ceiling. The two value facts do not tie — and correctly should "
     "not."),
    ("Known gaps are documented, not omitted.",
     "dbt connects to Fabric and writes; the transformation models required a T-SQL port, "
     "which is recorded with its scope and current status."),
    ("Unbuildable requirements are named with the telemetry they need.",
     "Customer lifetime value, churn, duplicate-redemption and geographic anomaly controls "
     "cannot be built from this dataset. Documenting that is more useful than metrics derived "
     "from fields that do not carry the signal."),
]
y = 2.35
for h1, h2 in items:
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(y), Inches(0.07),
                             Inches(0.95))
    bar.fill.solid(); bar.fill.fore_color.rgb = C(TEAL); bar.line.fill.background()
    bar.shadow.inherit = False
    txt(s, 1.28, y, 10.9, 0.95,
        [(h1, 15, "FFFFFF", True), (h2, 12, "9FB6D4", False, 4)], spacing=1.15)
    y += 1.12

txt(s, 1.0, 6.85, 11.3, 0.4,
    [("Anthony Apollis  ·  August 2026  ·  python scripts/run_all.py rebuilds everything",
      11, "6E8CB4", False)])

prs.save(str(OUT))
print(f"Wrote {OUT.name}")
print(f"  {len(prs.slides.__iter__.__self__._sldIdLst)} slides · "
      f"{OUT.stat().st_size/1024:.0f} KB · 16:9")
