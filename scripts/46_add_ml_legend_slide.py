"""
46_add_ml_legend_slide.py — a legend for the machine learning slide.

Slide 11 states five models and five metrics. Anyone who works with models reads it fine;
anyone who does not sees "AUC 0.620" and "silhouette 0.185" and has no way to tell whether
those are good, bad or meaningless. In a mixed panel that is most of the room.

This inserts a legend immediately after it, in two columns: what each METRIC means and what
each METHOD does, both in plain English, each with the number from the previous slide and a
one-line verdict. The point is to let the speaker say "0.620 out of a possible 0.621" and have
the audience already know what the scale is.

It could not go on slide 11 itself — that slide's last card already ends at 7.8 inches on a
7.5 inch canvas, so there is no room left.

The slide is appended and then MOVED into position, because python-pptx can only add at the
end; the move is done on the slide id list directly.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
DECK = ROOT / "report" / "Merchant_Voucher_Intelligence_Presentation.pptx"

NAVY, TEAL, AMBER = "12305B", "0E8B8B", "E8A317"
RED, GREEN, GREY, INK = "C0392B", "1E8449", "5A6672", "12203A"
BG, CARD, LINE = "F2F5FA", "FFFFFF", "DCE4EF"
C = lambda h: RGBColor.from_string(h.lstrip("#").upper())

# The legend must sit immediately AFTER the machine learning slide. It was originally a
# hardcoded index, which broke the moment another slide was inserted earlier in the deck —
# the legend silently moved in FRONT of the slide it explains. Find the anchor by title.
ANCHOR_TITLE = "Machine learning"

METRICS = [
    ("Area under the ROC curve (AUC)  ·  0.620", TEAL, "HIGHER is better  (0.5 = useless, 1.0 = perfect)",
     "Ranking quality. ROC = receiver operating characteristic. Read it as: the chance the model scores a real redeemer above a non-redeemer.",
     "0.620 against a measured ceiling of 0.621 for this data — 99.8% of the signal that "
     "exists here. Low because the data has no customer features, not because the model is weak."),
    ("Mean absolute percentage error (MAPE)  ·  4.69%", TEAL, "LOWER is better  (0% = perfect, no upper limit)",
     "Average forecast error as a percentage of the actual value.",
     "A 30-day sales forecast wrong by about 4.7% on average. 49.7% better than assuming "
     "tomorrow equals today."),
    ("Mean absolute error (MAE)  ·  9.8 hours", TEAL, "LOWER is better  (0 = perfect, measured in hours)",
     "Average error in the units being predicted — here, hours.",
     "Predicted ticket resolution time is out by about 10 hours on average, against a spread "
     "running to 74 hours at the 90th percentile."),
    ("Silhouette score  ·  0.185", AMBER, "HIGHER is better  (-1 to 1; above 0.5 = well separated)",
     "How cleanly separated the clusters are.",
     "Low, and stated rather than hidden. The merchants form a continuum, not four distinct "
     "types — so k was chosen on business grounds and the statistical cost is declared."),
    ("Naive baseline", GREY, "BEAT IT or the model is not worth deploying",
     "The simplest possible prediction: tomorrow equals today.",
     "Every model is reported against it. A model that cannot beat the naive baseline is not "
     "worth deploying, whatever its headline score."),
]

METHODS = [
    ("Isolation Forest", TEAL, "",
     "Unsupervised — finds rows unlike the rest without being told what to look for.",
     "Given no hints, it recovered all four documented patterns in the data."),
    ("Histogram-based gradient boosting", TEAL, "",
     "Supervised — many small decision trees, each correcting the last.",
     "Used for redemption propensity and resolution time. Handles missing values and mixed "
     "column types without manual preparation."),
    ("Holt-Winters exponential smoothing", TEAL, "",
     "Classical time series — level, trend and repeating seasonality.",
     "Chosen over a heavier model because seven months of daily data supports weekly "
     "seasonality and nothing more."),
    ("K-means clustering, k = 4", AMBER, "",
     "Groups merchants by similarity into a chosen number of clusters.",
     "k=4 matches how the business already talks about accounts, which matters more than a "
     "marginally better statistical score nobody can action."),
    ("Time-based split", GREEN, "",
     "Train on earlier data, test on later data. Never random.",
     "Train Jan-May, test Jun-Jul. A random split lets the model see the future and reports a "
     "score that will not survive production."),
]


def txt(slide, x, y, w, h, runs, spacing=1.05, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for item in runs:
        text, size, colour, bold = item[0], item[1], item[2], item[3]
        space_before = item[4] if len(item) > 4 else 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = spacing
        p.space_before = Pt(space_before)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = C(colour)
        r.font.name = "Segoe UI"
    return tb


def card(slide, x, y, w, h, fill=CARD, edge=LINE):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                               Inches(w), Inches(h))
    c.fill.solid()
    c.fill.fore_color.rgb = C(fill)
    c.line.color.rgb = C(edge)
    c.line.width = Pt(0.75)
    c.shadow.inherit = False
    try:
        c.adjustments[0] = 0.04
    except Exception:
        pass
    return c


def _drop_slides_titled(prs, prefix):
    """Remove every slide whose first text starts with `prefix`.

    Written as two passes on purpose. The first version walked the slide id list while
    removing from it and indexed the slide collection by position — so each removal shifted
    everything after it and the loop deleted the WRONG slides. Three slides were lost from the
    end of the deck before it was noticed. Collect first, delete second.
    """
    doomed = []
    for sid, slide in zip(list(prs.slides._sldIdLst), list(prs.slides)):
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip().startswith(prefix):
                doomed.append(sid)
                break
    for sid in doomed:
        prs.part.drop_rel(sid.rId)
        prs.slides._sldIdLst.remove(sid)
    return len(doomed)


prs = Presentation(str(DECK))
W, H = prs.slide_width, prs.slide_height
_n = _drop_slides_titled(prs, "What the numbers mean")
if _n:
    print(f"  removed {_n} previous copy/copies of this slide")


slide = prs.slides.add_slide(prs.slide_layouts[6])
bgr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
bgr.fill.solid()
bgr.fill.fore_color.rgb = C(BG)
bgr.line.fill.background()
bgr.shadow.inherit = False

band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(1.0))
band.fill.solid()
band.fill.fore_color.rgb = C(NAVY)
band.line.fill.background()
band.shadow.inherit = False
txt(slide, 0.55, 0.20, 11.0, 0.6, [("What the numbers mean", 26, "FFFFFF", True)])
txt(slide, 0.55, 0.66, 12.2, 0.3,
    [("A legend for the previous slide — every metric and method in plain English",
      12, "9FB6D4", False)])

COL_W, ROW_H, GAP = 6.05, 0.92, 0.06
for col, (heading, items, accent) in enumerate(
        [("Metrics — how good is good", METRICS, TEAL),
         ("Methods — what each model does", METHODS, NAVY)]):
    x = 0.55 + col * (COL_W + 0.15)
    txt(slide, x + 0.05, 1.22, COL_W, 0.3, [(heading, 13.5, accent, True)])
    y = 1.60
    for term, tcol, direction, plain, why in items:
        card(slide, x, y, COL_W, ROW_H)
        txt(slide, x + 0.18, y + 0.10, COL_W - 0.36, ROW_H - 0.2,
            [(term, 11, tcol, True),
             (direction, 9, RED if "LOWER" in direction else GREEN, True, 2),
             (plain, 9.5, INK, False, 2),
             (why, 8.5, GREY, False, 1)], spacing=1.02)
        y += ROW_H + GAP

foot_y = 1.60 + 5 * (ROW_H + GAP) + 0.10
card(slide, 0.55, foot_y, 12.25, 0.80, fill="EEF3F9", edge=NAVY)
txt(slide, 0.95, foot_y + 0.10, 11.5, 0.62,
    [("The four numbers that decide whether a model is any good — in this order",
      12, NAVY, True),
     ("1. Does it beat the naive baseline?  If not, stop — nothing else matters.      "
      "2. Train score vs test score.  A wide gap means it memorised rather than learned.      "
      "3. The headline metric against its CEILING, not against 1.0.      "
      "4. How was it validated?  A great score from a random split on time-series data is "
      "worthless.", 10, INK, False, 4)], spacing=1.15)

# python-pptx appends; move the new slide into position after AFTER_SLIDE.
anchor_pos = None
for _i, _s in enumerate(prs.slides, 1):
    for _sh in _s.shapes:
        if _sh.has_text_frame and _sh.text_frame.text.strip().startswith(ANCHOR_TITLE):
            anchor_pos = _i
            break
    if anchor_pos:
        break
if anchor_pos is None:
    raise SystemExit(f"cannot place the legend: no slide titled {ANCHOR_TITLE!r}")

sldIdLst = prs.slides._sldIdLst
_all = list(sldIdLst)
sldIdLst.remove(_all[-1])
sldIdLst.insert(anchor_pos, _all[-1])

slide.notes_slide.notes_text_frame.text = (
    "[45 sec] Only use this slide if the room needs it — if they are technical, skip straight "
    "past.\n\n"
    "The line that does the work: \"0.620 against a ceiling of 0.621 — so the model is "
    "capturing 99.8 percent of the signal that exists in this data. It looks low because "
    "there is no customer-level information in these files, not because the model is poor.\"\n\n"
    "The silhouette score is the honest one. 0.185 is low, and it is on the slide rather than "
    "hidden: the merchants form a continuum rather than four clean types, so k was chosen on "
    "business grounds and the statistical cost is declared.\n\n"
    "If asked why no deep learning: seven months of daily data on 25 merchants. The constraint "
    "is the data, not the algorithm."
)

prs.save(str(DECK))

chk = Presentation(str(DECK))
titles = []
for i, s in enumerate(chk.slides, 1):
    for sh in s.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            titles.append((i, sh.text_frame.text.strip().splitlines()[0][:44]))
            break
pos = [i for i, t in titles if t.startswith("What the numbers mean")]
notes = sum(1 for s in chk.slides
            if s.has_notes_slide and s.notes_slide.notes_text_frame.text.strip())

print(f"  legend inserted at slide {pos[0] if pos else '?'} of {len(chk.slides)}")
print(f"  {len(METRICS)} metrics and {len(METHODS)} methods explained")
print(f"  speaker notes on {notes} of {len(chk.slides)} slides")
for i, t in titles[9:13]:
    print(f"    {i:>2}. {t}")

if not pos or pos[0] != anchor_pos + 1:
    raise SystemExit(f"legend landed at {pos}, expected slide {anchor_pos + 1} "
                     f"(immediately after {ANCHOR_TITLE!r})")
