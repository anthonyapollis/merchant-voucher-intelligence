"""
47_add_dbt_value_slide.py — what dbt and the snapshot actually earned.

The deck showed the dbt lineage graph but never said what dbt BOUGHT. A DAG proves a tool was
used; it does not answer "why not just write Python". The Word report has that argument; the
deck did not, so a panel watching the walkthrough got the diagram without the case.

Framed as what would be LOST without each capability, because that is the form the argument
has to survive: a feature list invites "so what", a loss list does not.

The snapshot gets the largest card deliberately. It is the only item here that captures
something the source destroys — everything else could be rebuilt from the CSVs tomorrow.

Inserted after the lineage slide so the diagram is explained immediately after it is shown.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
DECK = ROOT / "report" / "Merchant_Voucher_Intelligence_Presentation.pptx"

NAVY, TEAL, AMBER = "12305B", "0E8B8B", "E8A317"
RED, GREEN, PURPLE, GREY, INK = "C0392B", "1E8449", "7B4B94", "5A6672", "12203A"
BG, CARD, LINE = "F2F5FA", "FFFFFF", "DCE4EF"
C = lambda h: RGBColor.from_string(h.lstrip("#").upper())

TITLE = "What dbt earned"
AFTER_SLIDE = 4           # straight after the lineage graph

EARNED = [
    ("15 relationships tests", TEAL,
     "Without them: a broken join returns FEWER ROWS, not an error. Totals quietly shrink and "
     "nobody is told."),
    ("22 unique · 56 not_null · 2 composite", TEAL,
     "Without them: duplicate-proofing is a claim rather than a check. A fact gaining a second "
     "row per key double-counts every measure built on it."),
    ("16 accepted_values · 11 expression · 5 range", TEAL,
     "Without them: a new Priority value or a negative sales figure reaches the report looking "
     "entirely legitimate."),
    ("5 seeds", TEAL,
     "Without them: the commercial team needs a code deployment to change a margin band."),
    ("21 models, documented from the manifest", TEAL,
     "Without them: the ERD and lineage are drawn by hand and start drifting from the code the "
     "day they are written."),
]


def txt(slide, x, y, w, h, runs, spacing=1.05, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for item in runs:
        text, size, colour, bold = item[0], item[1], item[2], item[3]
        sb = item[4] if len(item) > 4 else 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = spacing
        p.space_before = Pt(sb)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = C(colour)
        r.font.name = "Segoe UI"
    return tb


def card(slide, x, y, w, h, fill=CARD, edge=LINE):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                               Inches(w), Inches(h))
    c.fill.solid()
    c.fill.fore_color.rgb = C(fill)
    c.line.color.rgb = C(edge)
    c.line.width = Pt(0.75)
    c.shadow.inherit = False
    try:
        c.adjustments[0] = 0.04
    except Exception:
        pass
    return c


def _drop_slides_titled(prs, prefix):
    """Remove every slide whose first text starts with `prefix`.

    Written as two passes on purpose. The first version walked the slide id list while
    removing from it and indexed the slide collection by position — so each removal shifted
    everything after it and the loop deleted the WRONG slides. Three slides were lost from the
    end of the deck before it was noticed. Collect first, delete second.
    """
    doomed = []
    for sid, slide in zip(list(prs.slides._sldIdLst), list(prs.slides)):
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip().startswith(prefix):
                doomed.append(sid)
                break
    for sid in doomed:
        prs.part.drop_rel(sid.rId)
        prs.slides._sldIdLst.remove(sid)
    return len(doomed)


prs = Presentation(str(DECK))
W, H = prs.slide_width, prs.slide_height
_n = _drop_slides_titled(prs, TITLE)
if _n:
    print(f"  removed {_n} previous copy/copies of this slide")
H_IN = H / 914400


slide = prs.slides.add_slide(prs.slide_layouts[6])
bgr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
bgr.fill.solid(); bgr.fill.fore_color.rgb = C(BG); bgr.line.fill.background()
bgr.shadow.inherit = False

band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(1.0))
band.fill.solid(); band.fill.fore_color.rgb = C(NAVY); band.line.fill.background()
band.shadow.inherit = False
txt(slide, 0.55, 0.20, 11.0, 0.6, [(TITLE, 26, "FFFFFF", True)])
txt(slide, 0.55, 0.66, 12.2, 0.3,
    [("Stated as what would be LOST without it — a feature list invites “so what”",
      12, "9FB6D4", False)])

# ---- left: the contracts
txt(slide, 0.60, 1.22, 6.2, 0.3, [("Contracts that stop a bad build", 13.5, TEAL, True)])
y = 1.60
for term, tcol, loss in EARNED:
    card(slide, 0.55, y, 6.05, 0.90)
    txt(slide, 0.73, y + 0.11, 5.7, 0.72,
        [(term, 11.5, tcol, True), (loss, 9.5, GREY, False, 3)], spacing=1.05)
    y += 0.96

# ---- right: the snapshot, given the space it deserves
txt(slide, 6.80, 1.22, 6.0, 0.3, [("The snapshot — the one thing nothing else does",
                                   13.5, PURPLE, True)])
card(slide, 6.75, 1.60, 6.05, 2.70, fill="F6F2FA", edge=PURPLE)
txt(slide, 6.95, 1.78, 5.65, 2.4, [
    ("MerchantReference is a CURRENT-STATE extract.", 12, PURPLE, True),
    ("Every load overwrites active status, account manager, region and target. Two merchants "
     "read “At Risk” today — and without a snapshot, “when did that change, and did "
     "sales fall before or after?” is unanswerable forever.", 10.5, INK, False, 4),
    ("A Type 2 snapshot closes off the old version and opens a new one, so the history the "
     "source destroys is preserved in the warehouse.", 10.5, INK, False, 5),
    ("It uses the check strategy, not timestamp: the source has no reliable last-modified "
     "column — OnboardedDate is when the merchant joined, not when the row last changed.",
     9.5, GREY, False, 5),
], spacing=1.12)

card(slide, 6.75, 4.40, 6.05, 1.28, fill="EFF9F2", edge=GREEN)
txt(slide, 6.95, 4.56, 5.65, 1.05, [
    ("Proven, not asserted — 7/7 assertions", 12, GREEN, True),
    ("_test_scd2.py simulates a real change, re-runs the snapshot and checks: a new version is "
     "created, exactly one is current, the superseded row keeps the OLD value and is closed "
     "off, numbering stays contiguous, and unchanged merchants gain nothing. Then it restores "
     "state.", 9.5, INK, False, 4),
], spacing=1.10)

# ---- the honest limit
foot_y = 6.44
card(slide, 0.55, foot_y, 12.25, 0.80, fill="FEF8EC", edge=AMBER)
txt(slide, 0.95, foot_y + 0.11, 11.5, 0.62, [
    ("The honest limit of the argument", 11.5, AMBER, True),
    ("The transformation itself could have been written in Python — and it was. That "
     "duplication is deliberate: two implementations of the same logic must agree, and the "
     "reconciliation between them is what caught the R984,046 date-window defect that all 132 "
     "dbt tests passed straight through.", 10, INK, False, 3),
], spacing=1.12)

ids = prs.slides._sldIdLst
new = list(ids)[-1]
ids.remove(new)
ids.insert(AFTER_SLIDE, new)

slide.notes_slide.notes_text_frame.text = (
    "[90 sec] This is the answer to \"why not just write Python\".\n\n"
    "Left column: every one of those is a contract that stops a bad build. The line that "
    "lands is the first one — a broken join does not error, it returns fewer rows, so the "
    "totals quietly shrink and nobody is told.\n\n"
    "Right column is the one to dwell on. MerchantReference is a current-state extract; every "
    "load overwrites status, owner and target. Two merchants read At Risk today. Without the "
    "snapshot, when that changed is gone forever — and that is the question an account manager "
    "actually asks.\n\n"
    "Close on the amber card yourself, before they raise it: the transformation could have "
    "been Python, and it was. The duplication is the point — it is what caught the R984,046 "
    "defect that every dbt test passed straight through."
)

prs.save(str(DECK))

chk = Presentation(str(DECK))
h_in = chk.slide_height / 914400
over = []
pos = None
for i, s in enumerate(chk.slides, 1):
    mx = max((sh.top + sh.height) / 914400 for sh in s.shapes)
    if mx > h_in + 0.01:
        over.append((i, round(mx, 2)))
    for sh in s.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip().startswith(TITLE):
            pos = i
            break

notes = sum(1 for s in chk.slides
            if s.has_notes_slide and s.notes_slide.notes_text_frame.text.strip())
print(f"  '{TITLE}' inserted at slide {pos} of {len(chk.slides)}")
print(f"  speaker notes on {notes} of {len(chk.slides)} slides")
print(f"  overflowing slides: {over if over else 'none'}")

if pos != AFTER_SLIDE + 1:
    raise SystemExit(f"landed at slide {pos}, expected {AFTER_SLIDE + 1}")
if over:
    raise SystemExit(f"content runs off the canvas on {over}")
