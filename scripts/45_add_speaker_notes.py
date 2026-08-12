"""
45_add_speaker_notes.py — speaker notes for the walkthrough deck.

The deck was built to be TALKED THROUGH: one idea per slide, the number in large type, and the
reasoning small. That only works if the reasoning is somewhere the speaker can see and the
audience cannot. These are those notes.

They are written as things to SAY, not as a second copy of the slide. Each carries a rough
duration, the one number that has to land, and — where the slide invites a challenge — the
answer to the question it invites. Total runtime is about 22 minutes at a normal pace, which
leaves room in a typical 45-minute session for questions.

Register matters: the audience is a mixed panel that may include a CIO and a COO or CFO, so
the notes lead with the business consequence and keep the engineering underneath it, ready if
someone asks rather than volunteered.
"""
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
DECK = ROOT / "report" / "Merchant_Voucher_Intelligence_Presentation.pptx"

NOTES = {
1: """[30 sec] Open here — do not read the slide.

"This is the second-round practical task. Four CSV files came in. What I will walk you through
is what I built from them, and more to the point, how I know the numbers are right."

Set the frame early: this is about trustworthy reporting, not about a dashboard. If anyone in
the room owns finance, that is the thing they care about.""",

2: """[90 sec] The line that matters: every figure in this deck is reproducible — one command
rebuilds the whole thing in about two minutes.

Then the point of the slide: the gold layer is built TWICE, once in Python and once in SQL,
and they must agree to the cent before anything ships. That turns "the numbers tie" from a
claim into something checkable.

Plant the two defects it caught, but do not explain them — slide 13 does that. Just establish
that the checking found real things.""",

3: """[60 sec] Follow the row left to right. A file lands here, and by the time it reaches
Power BI it has passed a quality gate that stops the build if it fails.

The stage worth pausing on is the gate between silver and gold. If a test fails, gold is NOT
rebuilt, so the report keeps showing yesterday's correct numbers rather than today's wrong
ones. Stale and right beats fresh and wrong — that is a deliberate choice, not a limitation.""",

4: """[45 sec] This diagram is generated from dbt's manifest, not drawn by hand, so it cannot
drift from the code.

That is the whole point. Documentation maintained separately from code is wrong within a
month. Move on quickly unless asked — a technical audience will notice the lineage without you
labouring it.""",

5: """[60 sec] Be blunt about the starting point. Four flat files. The same merchant name in
four places, nothing enforcing they agree, every value arriving as text, and no calendar at
all.

The dashed lines are joins nobody validates — string matches that would fail silently. This
slide exists to make the next one mean something.""",

6: """[75 sec] Same data, after modelling. Fifteen foreign keys, and each one is a test that
runs on every build.

The line for a technical room: a broken join does not throw an error, it returns FEWER ROWS.
The totals quietly get smaller and nobody is told. That is why they are tested rather than
trusted.

If asked why dimensions are flat rather than fully normalised: at this size normalising saves
no meaningful storage and costs an extra join on every query. That is Kimball's own guidance,
not a shortcut.""",

7: """[75 sec] Expect a challenge, so raise it yourself. The README suggests five tables and I
built fourteen — that needs justifying table by table, not waving through.

Four come from supplied CSVs. One is DimDate, which the README asks for but does not supply.
Four exist because the brief requires them. One is a grain necessity. Four are mine, and they
are labelled as mine in the model itself.

The concession if pushed: those last four are the exposed surface. Say so before they do.""",

8: """[45 sec] Slow down and let the numbers sit. R65.5 million, 510,127 transactions, 84.2
percent redemption.

These are the same figures in the Power BI report, the Excel pack and the SQL — because they
come from one model, not three spreadsheets.

For a CFO, that sentence is the entire slide.""",

9: """[2 min] The brief's actual question set. The panel will care most about this slide. One
sentence each, then stop.

Two worth dwelling on. Voucher type: Airtime 92.8 percent against Gaming 76.0 — a 16.9 point
spread, and time-to-redeem is flat across types, so the difference is WHETHER customers
redeem, not how quickly.

And where management should focus: ranking by percentage decline sends you to the smallest
merchants. Ranking by RAND at risk sends you to the ones that matter.""",

10: """[2 min] The intellectual honesty slide. Take your time here.

The obvious analysis says tickets are associated with weaker performance — correlation minus
0.56. But it is confounded by size: bigger merchants have more of everything.

Control for size and it collapses to minus 0.20. So the portfolio-level rule is not supported.
The real signal is event-level — Umhlanga, tickets up 693 percent while sales fell 42.5
percent in the same month.

If they remember one slide about how you think, make it this one.""",

11: """[90 sec] Keep it short and resist the jargon.

Five models, all validated on time-based splits — never a random split, because a random split
on time-series data lets the model see the future and reports a score you cannot reproduce in
production.

The honest line: the redemption model reaches AUC 0.620 against a theoretical ceiling of 0.621
for this dataset. That is 99.8 percent of the available signal. It is low because the data has
no customer-level features, not because the model is weak.""",

12: """[45 sec] Five provinces, and the useful finding is where the business is NOT.

Eastern Cape is the only region below its own peak. Keep it brief — the map is doing the work,
and the detail sits on the Business Answers page if anyone wants it.""",

13: """[2 min] Your strongest slide with a finance audience.

Twenty-eight controls. Twenty-seven tie exactly. One is a documented rounding convention. And
one variance of R43.5 million is recorded as EXPECTED, with its reason attached.

Explain that last part: the sales fact and the voucher fact count different populations —
510,127 transactions against 120,969 vouchers. They must NOT agree. Recording it as expected
is what stops someone escalating it as a break every quarter.

Then the R984,046 story, about forty-five seconds. Lead with the mechanism, not the mistake:
two implementations disagreed, and the gap pointed at a date window no schema test could see.""",

14: """[90 sec] Frame it as prevent, detect, correct — detection alone is not a control
framework.

The part that lands: four controls CANNOT be built on this data, and each one names what it
would need. Duplicate redemption needs the raw event log including failed attempts. Customer
velocity needs a customer identifier, which does not exist in any of the four files.

Saying what you cannot do is more useful than a dashboard implying you can.""",

15: """[90 sec] Provisioned by REST API, authenticated with a CLI token — no service principal
secret stored anywhere in the repository.

The number to land: the medallion pipeline completed in five minutes fifty-four seconds, and
the gold tables in the Warehouse tie to R65,521,298.75 — the same figure as the local build,
to the cent.

The cost-control card is deliberate. Unattended compute does not stop on its own, so there is
a kill switch and a hard end date. In a room with a CIO that is worth ten seconds.""",

16: """[2 min] Ranked by revenue at risk, not severity of decline. Say that first.

Umhlanga: R571,518 at risk, sales down 42.5 percent and tickets up 693 percent in the same
month. Operational failure causing commercial damage — which means it is plausibly
recoverable, and that is why it is first.

Each recommendation names a merchant, a number and an owner. Account manager is carried onto
the risk register precisely so this resolves to a person rather than to a queue.""",

17: """[90 sec] Questions the solution raises beyond the brief. Use this to show you thought
past the task.

If time is short this is the slide to drop. It is fine to say there is more detail here if
useful, and move on.""",

18: """[60 sec] Close on principles, not features.

Figures are verifiable rather than asserted. The gold layer is built twice by independent
implementations. Assumptions are labelled as assumptions — of the 34 source columns, 12 are
documented by the supplier and 22 are inferred, and the documentation says which is which.

Then stop talking. Do not summarise the deck. Invite questions.""",
}


def main():
    if not DECK.exists():
        raise SystemExit(f"deck not found: {DECK}")

    prs = Presentation(str(DECK))
    n = len(prs.slides)

    extra = [i for i in NOTES if i > n]
    if extra:
        raise SystemExit(f"notes written for slides that do not exist: {extra}")

    words = 0
    for i, slide in enumerate(prs.slides, 1):
        text = NOTES.get(i)
        if not text:
            continue
        slide.notes_slide.notes_text_frame.text = text.strip()
        words += len(text.split())

    prs.save(str(DECK))

    # Re-open and confirm the notes actually persisted. python-pptx creates the notes slide
    # lazily, so a silent failure here would leave a deck that looks fine until it is opened
    # in front of a room.
    chk = Presentation(str(DECK))
    with_notes = [i for i, s in enumerate(chk.slides, 1)
                  if s.has_notes_slide and s.notes_slide.notes_text_frame.text.strip()]

    print(f"  speaker notes on {len(with_notes)} of {n} slides")
    print(f"  {words} words — roughly {words / 130:.0f} minutes of speaking, "
          f"leaving room for questions in a 45-minute session")
    print(f"  {DECK.name}  {DECK.stat().st_size / 1024:.0f} KB")
    print("  In PowerPoint: View > Notes Page, or drag up the pane below the slide.")

    if len(with_notes) != n:
        missing = [i for i in range(1, n + 1) if i not in with_notes]
        raise SystemExit(f"slides without speaker notes: {missing}")


main()
